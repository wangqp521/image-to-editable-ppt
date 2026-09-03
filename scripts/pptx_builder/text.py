"""Editable TextBox, Paragraph, Run, and native-list renderer."""

from __future__ import annotations

import re
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

from lib.capabilities import (
    ATOMIC_CAPABILITY_METADATA,
    CANONICAL_VALUES,
    TEXT_CONTRACT_ALLOWED_FIELDS,
    TEXT_RUN_BASELINE_MAX,
    TEXT_RUN_BASELINE_MIN,
    TEXT_RUN_MODERN_ALLOWED_FIELDS,
)
from lib.error_codes import ContractIssue
from lib.schema_contracts import (
    KIND_CONTENT_FIELDS,
    KIND_REQUIRED_CONTENT_FIELDS,
    KIND_REQUIRED_STYLE_FIELDS,
    KIND_STYLE_FIELDS,
)

from .common import RenderContext, register_renderer
from .ooxml import set_native_bullet
from .shapes import _apply_shadow, _validate_effects, apply_stroke, validate_stroke


_RGB = re.compile(r"#[0-9A-Fa-f]{6}")
_ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
    "distributed": PP_ALIGN.DISTRIBUTE,
}
_VERTICAL_ALIGNMENTS = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


def _issue(path: str, detail: str, capability: str | None = None) -> ContractIssue:
    return ContractIssue("UNSUPPORTED_CAPABILITY", path, detail, capability)


def _unknown_field_issue(
    value: dict[str, Any], allowed: frozenset[str], path: str
) -> ContractIssue | None:
    unknown = sorted(set(value) - allowed)
    if not unknown:
        return None
    field = unknown[0]
    return _issue(f"{path}.{field}", f"unknown typography field: {field}")


def _is_number(value: Any) -> bool:
    return type(value) in {int, float}


def _set_font_family(run: Any, family: str) -> None:
    run.font.name = family
    properties = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs", "a:sym"):
        node = properties.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            properties.append(node)
        node.set("typeface", family)


def _apply_run_style(run: Any, style: dict[str, Any], family: str) -> None:
    _set_font_family(run, family)
    font = run.font
    font.size = Pt(style["font_size"])
    font.bold = style["font_weight"] >= 600
    font.italic = style["italic"]
    font.underline = style["underline"]
    font.color.rgb = RGBColor.from_string(style["color"].removeprefix("#"))
    properties = run._r.get_or_add_rPr()
    properties.set("strike", "sngStrike" if style["strike"] else "noStrike")
    properties.set("baseline", str(style["baseline"]))
    properties.set("spc", str(round(style["letter_spacing"] * 100)))


def _add_paragraph_runs(
    paragraph: Any,
    contract: dict[str, Any],
    start: int,
    end: int,
) -> None:
    text = contract["text"]
    soft_breaks = {
        value
        for value in contract["text_box"].get("soft_breaks", [])
        if type(value) is int and start <= value <= end
    }
    boundaries = {start, end, *soft_breaks}
    for style in contract["runs"]:
        if style["start"] < end and style["end"] > start:
            boundaries.add(max(start, style["start"]))
            boundaries.add(min(end, style["end"]))
    ordered = sorted(boundaries)
    for index, left in enumerate(ordered[:-1]):
        if left in soft_breaks:
            paragraph.add_line_break()
        right = ordered[index + 1]
        if right <= left:
            continue
        style = next(
            item
            for item in contract["runs"]
            if item["start"] <= left < item["end"]
        )
        run = paragraph.add_run()
        run.text = text[left:right]
        _apply_run_style(run, style, contract["selected_font"])
    if end in soft_breaks:
        paragraph.add_line_break()


def _apply_paragraph_contract(
    paragraph: Any,
    paragraph_contract: dict[str, Any],
    contract: dict[str, Any],
    index: int,
    *,
    slide_part: Any | None = None,
) -> None:
    path = f"modules.typography.items.{contract['element_id']}.paragraphs[{index}]"
    paragraph.alignment = _ALIGNMENTS[paragraph_contract["alignment"]]
    paragraph.line_spacing = paragraph_contract["line_spacing"]
    paragraph.space_before = Pt(paragraph_contract["space_before"])
    paragraph.space_after = Pt(paragraph_contract["space_after"])
    properties = paragraph._p.get_or_add_pPr()
    if "margin_left" in paragraph_contract:
        properties.set("marL", str(round(paragraph_contract["margin_left"])))
    if "indent" in paragraph_contract:
        properties.set("indent", str(round(paragraph_contract["indent"])))
    list_contract = paragraph_contract["list"]
    if list_contract["is_list"]:
        bullet_contract = dict(list_contract)
        bullet_contract["indent"] = round(paragraph_contract["indent"])
        set_native_bullet(
            paragraph,
            bullet_contract,
            f"{path}.list",
            slide_part=slide_part,
        )
    _add_paragraph_runs(
        paragraph,
        contract,
        paragraph_contract["start"],
        paragraph_contract["end"],
    )


def apply_text_contract(
    text_frame: Any, contract: dict[str, Any], *, slide_part: Any | None = None
) -> None:
    """Apply a validated typography item without inferring missing style."""
    text_box = contract["text_box"]
    margins = text_box["margins"]
    text_frame.clear()
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = Emu(round(margins["left"]))
    text_frame.margin_right = Emu(round(margins["right"]))
    text_frame.margin_top = Emu(round(margins["top"]))
    text_frame.margin_bottom = Emu(round(margins["bottom"]))
    text_frame.word_wrap = text_box["wrap"]
    text_frame.vertical_anchor = _VERTICAL_ALIGNMENTS[text_box["vertical_alignment"]]
    body_properties = text_frame._txBody.bodyPr
    body_properties.set("vertOverflow", "overflow")
    body_properties.set("horzOverflow", "overflow")
    for index, paragraph_contract in enumerate(contract["paragraphs"]):
        paragraph = (
            text_frame.paragraphs[0]
            if index == 0
            else text_frame.add_paragraph()
        )
        _apply_paragraph_contract(
            paragraph,
            paragraph_contract,
            contract,
            index,
            slide_part=slide_part,
        )


def render(context: RenderContext, element: dict[str, Any]) -> None:
    """Render one text element from the pre-indexed typography contract."""
    contract = context.typography[element["element_id"]]
    x, y, width, height = element["slide_bbox"]
    shape = context.slide.shapes.add_textbox(
        Emu(x), Emu(y), Emu(width), Emu(height)
    )
    if element["style"].get("fill") == "noFill":
        shape.fill.background()
    if "line" in element["style"]:
        apply_stroke(shape, element["style"]["line"])
    else:
        shape.line.fill.background()
    _apply_shadow(shape, element["style"].get("effects"))
    shape.rotation = element["style"].get("rotation", 0)
    apply_text_contract(shape.text_frame, contract, slide_part=context.slide.part)
    context.registry.register(
        element["element_id"],
        shape,
        "sp",
        semantic_kind="text",
        selected_mode=context.representation_modes[element["element_id"]],
        text_summary=contract["text"],
        font_declarations=(contract["selected_font"],),
    )


class TextRenderer:
    kind = "text"
    supported_fields = KIND_STYLE_FIELDS["text"] | KIND_CONTENT_FIELDS["text"]
    supported_values = {"bullet_type": CANONICAL_VALUES["bullet_type"]}
    required_fields = (
        KIND_REQUIRED_STYLE_FIELDS["text"]
        | KIND_REQUIRED_CONTENT_FIELDS["text"]
    )
    capability_ids = frozenset(
        capability_id
        for capability_id, field in ATOMIC_CAPABILITY_METADATA.items()
        if field == "text"
    ) | frozenset(
        f"bullet_type.{value}" for value in CANONICAL_VALUES["bullet_type"]
    )

    def validate_contract(
        self, element: dict[str, Any], context: RenderContext
    ) -> list[ContractIssue]:
        path = f"elements.{element.get('element_id', '<unknown>')}"
        style = element.get("style", {})
        content = element.get("content", {})
        fields = set(style) | set(content)
        unsupported = sorted(fields - self.supported_fields)
        if unsupported:
            return [_issue(path, f"unsupported text fields: {', '.join(unsupported)}")]
        if not self.required_fields <= fields:
            return [_issue(f"{path}.content.text", "text is required")]
        if style.get("fill") not in {None, "noFill"}:
            return [_issue(f"{path}.style.fill", "only noFill is supported")]
        if "line" in style:
            issues = validate_stroke(
                style["line"], f"{path}.style.line", "shape.line"
            )
            if issues:
                return issues
        if "effects" in style:
            issues = _validate_effects(style["effects"], f"{path}.style.effects")
            if issues:
                return issues
        if "margins" in style:
            margins = style["margins"]
            expected = {"left", "right", "top", "bottom"}
            if (
                not isinstance(margins, dict)
                or set(margins) != expected
                or any(
                    not _is_number(margins[side]) or margins[side] < 0
                    for side in expected
                )
            ):
                return [_issue(
                    f"{path}.style.margins",
                    "all four text style margins must be non-negative numbers",
                    "text.frame.margins",
                )]
        if "vertical_alignment" in style and style["vertical_alignment"] not in _VERTICAL_ALIGNMENTS:
            return [_issue(
                f"{path}.style.vertical_alignment",
                "unsupported text style vertical alignment",
                "text.frame.vertical_alignment",
            )]
        if "wrap" in style and type(style["wrap"]) is not bool:
            return [_issue(
                f"{path}.style.wrap",
                "text style wrap must be boolean",
                "text.frame.wrap",
            )]
        rotation = style.get("rotation", 0)
        if not _is_number(rotation) or not -360 <= rotation <= 360:
            return [_issue(
                f"{path}.style.rotation",
                "rotation must be from -360 to 360 degrees",
            )]
        element_id = element.get("element_id")
        contract = context.typography.get(element_id)
        if not isinstance(contract, dict):
            return [_issue("modules.typography.items", "text element lacks typography")]
        contract_path = f"modules.typography.items.{element_id}"
        unknown_issue = _unknown_field_issue(
            contract, TEXT_CONTRACT_ALLOWED_FIELDS["item"], contract_path
        )
        if unknown_issue is not None:
            return [unknown_issue]
        if contract.get("text") != content.get("text"):
            return [_issue(f"{path}.content.text", "typography text must match element text")]
        if context.representation_modes.get(element_id) != "native":
            return [_issue(f"{path}.representation", "text renderer requires native mode")]
        if not isinstance(contract.get("selected_font"), str) or not contract["selected_font"]:
            return [_issue("modules.typography.items", "selected_font is required")]
        text_box = contract.get("text_box")
        if not isinstance(text_box, dict):
            return [_issue("modules.typography.items", "text_box is required")]
        unknown_issue = _unknown_field_issue(
            text_box,
            TEXT_CONTRACT_ALLOWED_FIELDS["text_box"],
            f"{contract_path}.text_box",
        )
        if unknown_issue is not None:
            return [unknown_issue]
        margins = text_box.get("margins")
        if (
            not isinstance(margins, dict)
            or any(
                not _is_number(margins.get(side))
                for side in ("left", "right", "top", "bottom")
            )
        ):
            return [_issue(
                f"modules.typography.items.{element_id}.text_box.margins",
                "all four TextBox margins must be numeric EMU",
                "text.frame.margins",
            )]
        if type(text_box.get("wrap")) is not bool:
            return [_issue(
                f"modules.typography.items.{element_id}.text_box.wrap",
                "TextBox wrap must be boolean",
                "text.frame.wrap",
            )]
        if text_box.get("overflow") is not True:
            return [_issue(
                f"modules.typography.items.{element_id}.text_box.overflow",
                "ordinary typography TextBox overflow must be true",
            )]
        if text_box.get("vertical_alignment") not in _VERTICAL_ALIGNMENTS:
            return [_issue(
                f"modules.typography.items.{element_id}.text_box.vertical_alignment",
                "unsupported TextBox vertical alignment",
                "text.frame.vertical_alignment",
            )]
        for style_field, contract_field in (
            ("margins", "margins"),
            ("vertical_alignment", "vertical_alignment"),
            ("wrap", "wrap"),
        ):
            if (
                style_field in style
                and style[style_field] != text_box.get(contract_field)
            ):
                return [_issue(
                    f"{path}.style.{style_field}",
                    f"text style {style_field} must match typography text_box",
                )]
        for run_index, run_contract in enumerate(contract.get("runs", [])):
            run_path = f"modules.typography.items.{element_id}.runs[{run_index}]"
            if not isinstance(run_contract, dict):
                return [_issue(run_path, "run must be an object")]
            unknown_issue = _unknown_field_issue(
                run_contract, TEXT_RUN_MODERN_ALLOWED_FIELDS, run_path
            )
            if unknown_issue is not None:
                return [unknown_issue]
            missing = sorted(TEXT_RUN_MODERN_ALLOWED_FIELDS - set(run_contract))
            if missing:
                return [_issue(
                    run_path,
                    f"missing run fields: {', '.join(missing)}",
                )]
            if any(
                type(run_contract[field]) is not bool
                for field in ("italic", "underline", "strike")
            ):
                return [_issue(run_path, "run style flags must be boolean")]
            baseline = run_contract["baseline"]
            if (
                type(baseline) is not int
                or not TEXT_RUN_BASELINE_MIN <= baseline <= TEXT_RUN_BASELINE_MAX
            ):
                return [_issue(
                    f"{run_path}.baseline",
                    "baseline must be an integer from -100000 to 100000",
                    "text.run.baseline",
                )]
            color = run_contract.get("color")
            if not isinstance(color, str) or _RGB.fullmatch(color) is None:
                return [_issue(
                    f"{run_path}.color",
                    "run color must be #RRGGBB",
                    "text.run.color",
                )]
        for paragraph_index, paragraph in enumerate(contract.get("paragraphs", [])):
            if not isinstance(paragraph, dict):
                return [_issue(
                    f"modules.typography.items.{element_id}.paragraphs[{paragraph_index}]",
                    "paragraph must be an object",
                )]
            paragraph_path = (
                f"modules.typography.items.{element_id}.paragraphs[{paragraph_index}]"
            )
            unknown_issue = _unknown_field_issue(
                paragraph,
                TEXT_CONTRACT_ALLOWED_FIELDS["paragraph"],
                paragraph_path,
            )
            if unknown_issue is not None:
                return [unknown_issue]
            if paragraph.get("alignment") not in _ALIGNMENTS:
                return [_issue(
                    f"modules.typography.items.{element_id}.paragraphs[{paragraph_index}].alignment",
                    "unsupported paragraph alignment",
                )]
            list_contract = paragraph.get("list")
            if isinstance(list_contract, dict):
                unknown_issue = _unknown_field_issue(
                    list_contract,
                    TEXT_CONTRACT_ALLOWED_FIELDS["list"],
                    f"{paragraph_path}.list",
                )
                if unknown_issue is not None:
                    return [unknown_issue]
            if isinstance(list_contract, dict) and list_contract.get("is_list") is True:
                bullet_type = list_contract.get("bullet_type")
                if bullet_type not in CANONICAL_VALUES["bullet_type"]:
                    return [_issue(
                        f"modules.typography.items.{element_id}.paragraphs[{paragraph_index}].list.bullet_type",
                        f"unsupported native bullet type: {bullet_type}",
                        f"bullet_type.{bullet_type}",
                    )]
        return []

    def render(self, element: dict[str, Any], context: RenderContext) -> None:
        render(context, element)


TEXT_RENDERER = TextRenderer()
register_renderer("text", TEXT_RENDERER)
