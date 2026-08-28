#!/usr/bin/env python3
"""Compile one passing schema v2 page specification into an editable PPTX."""

from __future__ import annotations

import argparse
import importlib.metadata
import json
import platform
import tempfile
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation

from lib.atomic_write import (
    atomic_write_bytes,
    atomic_write_json,
    publish_pair_no_overwrite,
    publish_pair_replace_current,
)
from lib.background_contracts import (
    resolved_element_mode_map,
    validate_background_prebuild,
)
from lib.capabilities import capability_manifest_sha256
from lib.error_codes import ContractIssue, ToolError
from lib.font_runtime import validate_font_runtime
from lib.hashing import canonical_json_sha256, file_sha256
from lib.path_contracts import find_user_controlled_symlink
from lib.representation_contracts import representation_summary, validate_representation_plan
from lib.schema_contracts import schema_envelope_issues
from lib.schema_io import index_elements, load_schema_v2
from lib.spec_identity import content_spec_sha256, input_spec_sha256
from normalize_native_list_ooxml import NormalizeError, normalize_pptx
from pptx_builder import (
    RENDERERS,
    ObjectRegistry,
    RenderContext,
    validate_renderer_contracts,
)


def _load_json_object(path: Path) -> dict[str, Any]:
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, UnicodeError, json.JSONDecodeError) as exc:
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE", str(path), "invalid prebuild report"
        ) from exc
    if not isinstance(value, dict):
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE",
            str(path),
            "prebuild report must be an object",
        )
    return value


def _freeze_spec(spec: dict[str, Any]) -> tuple[dict[str, Any], bytes]:
    """Return one immutable-by-convention canonical JSON snapshot and its bytes."""
    try:
        encoded = (
            json.dumps(
                spec,
                allow_nan=False,
                ensure_ascii=False,
                separators=(",", ":"),
                sort_keys=True,
            )
            + "\n"
        ).encode("utf-8")
        frozen = json.loads(encoded)
    except (TypeError, UnicodeError, ValueError, OverflowError, RecursionError) as exc:
        raise ToolError(
            "SPEC_IDENTITY_INVALID",
            "$",
            "spec must contain only finite JSON values",
        ) from exc
    return frozen, encoded


def compiler_identity() -> dict[str, str]:
    """Return the sorted content identity of all compiler runtime modules."""
    scripts_root = Path(__file__).resolve().parent
    paths = [scripts_root / "build_pptx_from_spec.py"]
    paths.extend((scripts_root / "lib").glob("*.py"))
    paths.extend((scripts_root / "pptx_builder").glob("*.py"))
    relative_paths = sorted(
        path.relative_to(scripts_root).as_posix()
        for path in paths
        if path.is_file()
    )
    return {
        relative_path: file_sha256(scripts_root / relative_path)
        for relative_path in relative_paths
    }


def compiler_sha256() -> str:
    """Return a traversal- and metadata-independent compiler digest."""
    return canonical_json_sha256(compiler_identity())


def _environment() -> dict[str, str]:
    return {
        "Pillow": importlib.metadata.version("Pillow"),
        "python": platform.python_version(),
        "python-pptx": importlib.metadata.version("python-pptx"),
    }


def _asset_fallbacks(spec: dict[str, Any]) -> list[dict[str, Any]]:
    items = spec["modules"]["representation_plan"]["items"]
    return [
        dict(item)
        for item in sorted(items, key=lambda value: value["source_fact_id"])
        if item["render_mode"] == "picture_asset"
    ]


def _background_summary(spec: dict[str, Any]) -> dict[str, int]:
    summary = {"native": 0, "background_picture": 0}
    for item in spec["modules"]["background"]["items"]:
        summary[item["selected_mode"]] += 1
    return summary


def _background_pictures(
    spec: dict[str, Any], element_report: dict[str, Any]
) -> list[dict[str, Any]]:
    pictures: list[dict[str, Any]] = []
    for item in sorted(
        spec["modules"]["background"]["items"],
        key=lambda value: value["background_id"],
    ):
        if item["selected_mode"] != "background_picture":
            continue
        element_id = item["bound_element_id"]
        objects = element_report[element_id]["objects"]
        media_hashes = {
            obj["media_sha256"]
            for obj in objects
            if obj["media_sha256"] is not None
        }
        if len(media_hashes) != 1:
            raise ToolError(
                "BUILD_OUTPUT_INCOMPLETE",
                f"elements.{element_id}",
                "background picture must register exactly one media hash",
            )
        pictures.append({**item, "media_sha256": next(iter(media_hashes))})
    return pictures


def _validate_output_paths(output_pptx: Path, build_report_path: Path) -> None:
    if output_pptx == build_report_path:
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE",
            str(output_pptx),
            "PPTX and build report paths must be distinct",
        )
    if (
        output_pptx in build_report_path.parents
        or build_report_path in output_pptx.parents
    ):
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE",
            str(build_report_path),
            "PPTX and build report paths must not contain one another",
        )


def _validate_raw_output_entry(destination: Path) -> None:
    """Reject symlink output entries before canonical path comparison."""
    try:
        symlink = find_user_controlled_symlink(destination)
        if symlink == destination:
            raise ToolError(
                "BUILD_OUTPUT_INCOMPLETE",
                str(destination),
                "output path already exists",
            )
        if symlink is not None:
            raise ToolError(
                "BUILD_OUTPUT_INCOMPLETE",
                str(symlink),
                "output parent must be an existing real directory",
            )
    except ToolError:
        raise
    except (OSError, RuntimeError, ValueError) as exc:
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE",
            str(destination),
            "cannot inspect output path",
        ) from exc


def _build_report(
    spec: dict[str, Any],
    schema_sha256: str,
    content_identity: str,
    input_identity: str,
    candidate: Path,
    element_report: dict[str, Any],
    normalization: dict[str, Any],
    build_report_path: Path,
    preferred_font: str,
    font_runtime: dict[str, Any] | None,
    runtime_preflight: dict[str, str] | None,
) -> dict[str, Any]:
    try:
        report = {
            "valid": True,
            "schema_version": 1,
            "schema_sha256": schema_sha256,
            "content_spec_sha256": content_identity,
            "input_spec_sha256": input_identity,
            "preferred_font": preferred_font,
            "runtime_preflight": (
                dict(runtime_preflight) if runtime_preflight is not None else None
            ),
            "font_runtime": dict(font_runtime) if font_runtime is not None else None,
            "compiler_sha256": compiler_sha256(),
            "capability_manifest_sha256": capability_manifest_sha256(),
            "pptx_sha256": file_sha256(candidate),
            "environment": _environment(),
            "elements": element_report,
            "representation_summary": representation_summary(spec),
            "asset_fallbacks": _asset_fallbacks(spec),
            "background_summary": _background_summary(spec),
            "background_pictures": _background_pictures(spec, element_report),
            "normalization": normalization,
            "warnings": [],
            "unsupported": [],
        }
        json.dumps(report, allow_nan=False, ensure_ascii=False, sort_keys=True)
    except ToolError:
        raise
    except (
        OSError,
        UnicodeError,
        TypeError,
        ValueError,
        OverflowError,
        KeyError,
        importlib.metadata.PackageNotFoundError,
    ) as exc:
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE",
            str(build_report_path),
            "cannot construct deterministic build report",
        ) from exc
    return report


def _ensure_output_parent(destination: Path) -> None:
    try:
        destination.parent.mkdir(parents=True, exist_ok=True)
    except OSError as exc:
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE",
            str(destination.parent),
            "cannot create output parent directory",
        ) from exc


def _raise_issue(issue: ContractIssue) -> None:
    raise ToolError(issue.code, issue.path, issue.detail, issue.capability)


def _typography_index(spec: dict[str, Any]) -> dict[str, dict[str, Any]]:
    module = spec.get("modules", {}).get("typography")
    items = module.get("items") if isinstance(module, dict) else None
    if not isinstance(items, list):
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE",
            "modules.typography.items",
            "typography items must be an array",
        )
    indexed: dict[str, dict[str, Any]] = {}
    for index, item in enumerate(items):
        element_id = item.get("element_id") if isinstance(item, dict) else None
        if not isinstance(element_id, str) or not element_id or element_id in indexed:
            raise ToolError(
                "BUILD_OUTPUT_INCOMPLETE",
                f"modules.typography.items[{index}].element_id",
                "typography element_id must be unique and non-empty",
            )
        indexed[element_id] = item
    return indexed


def _has_native_list(typography: dict[str, dict[str, Any]]) -> bool:
    return any(
        isinstance(paragraph, dict)
        and isinstance(paragraph.get("list"), dict)
        and paragraph["list"].get("is_list") is True
        for contract in typography.values()
        for paragraph in contract.get("paragraphs", [])
    )


def _validate_candidate(path: Path) -> None:
    try:
        with zipfile.ZipFile(path, "r") as archive:
            bad_member = archive.testzip()
            names = set(archive.namelist())
        if bad_member is not None or "ppt/presentation.xml" not in names:
            raise ValueError("invalid PPTX ZIP structure")
        Presentation(path)
    except (
        OSError,
        RuntimeError,
        ValueError,
        zipfile.BadZipFile,
        zipfile.LargeZipFile,
    ) as exc:
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE", str(path), "saved PPTX cannot be reopened"
        ) from exc


def _gate_contracts(
    spec: dict[str, Any],
    elements: dict[str, dict[str, Any]],
    typography: dict[str, dict[str, Any]],
) -> dict[str, str]:
    representation_issues = validate_representation_plan(spec)
    if representation_issues:
        _raise_issue(representation_issues[0])
    background_issues = validate_background_prebuild(spec)
    if background_issues:
        _raise_issue(background_issues[0])
    modes = resolved_element_mode_map(spec)
    renderer_issues = validate_renderer_contracts(
        spec, elements, modes, typography
    )
    if renderer_issues:
        _raise_issue(renderer_issues[0])
    return modes


def _gate_schema_envelopes(spec: dict[str, Any]) -> None:
    for path, detail in schema_envelope_issues(spec):
        raise ToolError("UNSUPPORTED_CAPABILITY", path, detail)


def compile_single_page(
    spec_path: str | Path,
    prebuild_report_path: str | Path,
    output_pptx: str | Path,
    build_report_path: str | Path,
    *,
    replace_current: bool = False,
) -> dict[str, Any]:
    """Compile and atomically publish one schema-bound PPTX/report pair."""
    spec_path = Path(spec_path).expanduser().resolve()
    prebuild_report_path = Path(prebuild_report_path).expanduser().resolve()
    output_pptx = Path(output_pptx).expanduser()
    build_report_path = Path(build_report_path).expanduser()
    _validate_raw_output_entry(output_pptx)
    _validate_raw_output_entry(build_report_path)
    output_pptx = output_pptx.resolve()
    build_report_path = build_report_path.resolve()
    _validate_output_paths(output_pptx, build_report_path)
    if not replace_current and (output_pptx.exists() or build_report_path.exists()):
        existing = output_pptx if output_pptx.exists() else build_report_path
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE", str(existing), "output path already exists"
        )
    if replace_current and output_pptx.exists() != build_report_path.exists():
        existing = output_pptx if output_pptx.exists() else build_report_path
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE",
            str(existing),
            "current PPTX and build report must either both exist or both be absent",
        )

    spec, spec_snapshot = _freeze_spec(load_schema_v2(spec_path))
    prebuild_report = _load_json_object(prebuild_report_path)
    schema_sha256 = canonical_json_sha256(spec)
    content_identity = content_spec_sha256(spec)
    input_identity = input_spec_sha256(spec)
    if prebuild_report.get("spec_sha256") != schema_sha256:
        raise ToolError(
            "SPEC_HASH_MISMATCH",
            "prebuild_report.spec_sha256",
            "prebuild report does not bind the current schema",
        )
    if (
        prebuild_report.get("valid") is not True
        or prebuild_report.get("stage") != "prebuild"
        or prebuild_report.get("errors") != []
    ):
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE",
            str(prebuild_report_path),
            "compiler requires a passing prebuild report",
        )
    snapshot = prebuild_report.get("snapshot")
    if not isinstance(snapshot, dict):
        raise ToolError(
            "PREBUILD_SNAPSHOT_MISMATCH",
            "prebuild_report.snapshot",
            "snapshot identity is required",
        )

    preferred_font = prebuild_report.get("preferred_font")
    if not isinstance(preferred_font, str) or not preferred_font.strip():
        raise ToolError(
            "BUILD_OUTPUT_INCOMPLETE",
            "prebuild_report.preferred_font",
            "a non-empty preferred font is required",
        )
    preferred_font = preferred_font.strip()
    font_runtime: dict[str, Any] | None = None
    runtime_preflight = prebuild_report.get("runtime_preflight")
    font_runtime_value = prebuild_report.get("font_runtime")
    if font_runtime_value is not None or runtime_preflight is not None:
        try:
            font_runtime = validate_font_runtime(font_runtime_value)
        except ValueError as exc:
            raise ToolError(
                "BUILD_OUTPUT_INCOMPLETE",
                "prebuild_report.font_runtime",
                str(exc),
            ) from exc
        if (
            not isinstance(runtime_preflight, dict)
            or not isinstance(runtime_preflight.get("path"), str)
            or not Path(runtime_preflight["path"]).is_absolute()
            or not isinstance(runtime_preflight.get("sha256"), str)
            or len(runtime_preflight["sha256"]) != 64
        ):
            raise ToolError(
                "BUILD_OUTPUT_INCOMPLETE",
                "prebuild_report.runtime_preflight",
                "runtime and font identities must be supplied together",
            )
        if font_runtime["family"] != preferred_font:
            raise ToolError(
                "BUILD_OUTPUT_INCOMPLETE",
                "prebuild_report.preferred_font",
                "preferred font does not match the optional strict runtime",
            )
    if (
        snapshot.get("path") != str(spec_path)
        or snapshot.get("sha256") != file_sha256(spec_path)
    ):
        raise ToolError(
            "PREBUILD_SNAPSHOT_MISMATCH",
            "prebuild_report.snapshot",
            "prebuild report does not bind the supplied build snapshot",
        )

    _gate_schema_envelopes(spec)
    elements = index_elements(spec)
    typography = _typography_index(spec)
    modes = _gate_contracts(spec, elements, typography)

    reading_order = {
        element_id: index for index, element_id in enumerate(spec["reading_order"])
    }
    ordered_elements = sorted(
        elements.values(),
        key=lambda element: (
            element["layer"],
            reading_order[element["element_id"]],
        ),
    )

    with tempfile.TemporaryDirectory(prefix="ia-pptx-compiler-") as directory:
        transaction = Path(directory)
        raw_candidate = transaction / "raw.pptx"
        normalized_candidate = transaction / "normalized.pptx"
        normalization_report_path = transaction / "normalization.json"
        spec_snapshot_path = transaction / "spec-snapshot.json"
        atomic_write_bytes(spec_snapshot_path, spec_snapshot)
        presentation = Presentation()
        presentation.slide_width, presentation.slide_height = spec["canvas"][
            "slide_size_emu"
        ]
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        registry = ObjectRegistry()
        context = RenderContext(
            slide=slide,
            spec=spec,
            representation_modes=modes,
            typography=typography,
            registry=registry,
            font_family=preferred_font,
        )
        for element in ordered_elements:
            RENDERERS[element["kind"]].render(element, context)
        element_report = registry.finalize(spec, modes)
        presentation.save(raw_candidate)

        if _has_native_list(typography):
            try:
                normalization_result = normalize_pptx(
                    raw_candidate,
                    spec_snapshot_path,
                    normalized_candidate,
                    normalization_report_path,
                )
                normalization = {
                    "applied": True,
                    "valid": normalization_result["valid"],
                    "paragraphs_checked": normalization_result[
                        "paragraphs_checked"
                    ],
                    "paragraphs_changed": normalization_result[
                        "paragraphs_changed"
                    ],
                }
            except NormalizeError as exc:
                raise ToolError(exc.code, exc.path, exc.detail) from exc
            except ToolError:
                raise
            except (
                OSError,
                ValueError,
                KeyError,
                TypeError,
                zipfile.BadZipFile,
                zipfile.LargeZipFile,
            ) as exc:
                raise ToolError(
                    "BUILD_OUTPUT_INCOMPLETE",
                    str(raw_candidate),
                    "native list normalization failed",
                ) from exc
            raw_candidate.unlink(missing_ok=True)
            candidate = normalized_candidate
        else:
            candidate = raw_candidate
            normalization = {"applied": False}

        _validate_candidate(candidate)
        report = _build_report(
            spec,
            schema_sha256,
            content_identity,
            input_identity,
            candidate,
            element_report,
            normalization,
            build_report_path,
            preferred_font,
            font_runtime,
            runtime_preflight,
        )
        _ensure_output_parent(output_pptx)
        _ensure_output_parent(build_report_path)
        report_candidate = transaction / "build-report.json"
        try:
            atomic_write_json(report_candidate, report)
            _load_json_object(report_candidate)
        except ToolError:
            raise
        except (OSError, TypeError, ValueError) as exc:
            raise ToolError(
                "BUILD_OUTPUT_INCOMPLETE",
                str(build_report_path),
                "cannot write build report candidate",
            ) from exc
        publisher = (
            publish_pair_replace_current if replace_current else publish_pair_no_overwrite
        )
        publisher(candidate, report_candidate, output_pptx, build_report_path)
    return report


def _parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--spec", type=Path, required=True)
    parser.add_argument("--prebuild-report", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--build-report", type=Path, required=True)
    parser.add_argument(
        "--replace-current",
        action="store_true",
        help="atomically replace the current PPTX/build-report pair if it exists",
    )
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = _parse_args(argv)
    try:
        report = compile_single_page(
            args.spec,
            args.prebuild_report,
            args.output,
            args.build_report,
            replace_current=args.replace_current,
        )
    except ToolError as exc:
        print(
            json.dumps(
                {"valid": False, "errors": [exc.as_dict()]}, ensure_ascii=False
            )
        )
        return 2
    print(json.dumps(report, ensure_ascii=False, sort_keys=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
