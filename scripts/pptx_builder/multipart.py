"""Matrix and status renderers for deterministic editable parts."""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

from lib.element_contracts import expand_multipart_parts, validate_element_contract
from lib.error_codes import ContractIssue, ToolError
from lib.geometry import (
    bbox_union,
    valid_font_size_pt,
    valid_nonnegative_coordinate32,
    validate_bbox,
)

from .common import RenderContext, register_renderer
from .ooxml import set_preset_shape_adjustments, set_shape_flips
from .shapes import (
    SHAPE_RENDERER,
    _SHAPE_TYPES,
    _apply_fill,
    _apply_shadow,
    apply_stroke,
)


_ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}
_VERTICAL_ALIGNMENTS = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}
_TEXT_STYLE_FIELDS = frozenset(
    {
        "font_name",
        "font_size",
        "font_weight",
        "color",
        "italic",
        "alignment",
        "vertical_alignment",
        "margins",
        "wrap",
    }
)
_MARGIN_FIELDS = frozenset({"left", "right", "top", "bottom"})


def _issue(path: str, detail: str) -> ContractIssue:
    return ContractIssue("UNSUPPORTED_CAPABILITY", path, detail)


def _normalized_fill(value: Any) -> Any:
    if isinstance(value, str) and value.startswith("#"):
        return {"type": "solid", "color": value, "opacity": 1.0}
    return value


def _part_shape_style(style: dict[str, Any]) -> dict[str, Any]:
    result = {key: value for key, value in style.items() if key != "text_style"}
    if "fill" in result:
        result["fill"] = _normalized_fill(result["fill"])
    return result


def _valid_rgb(value: Any) -> bool:
    return (
        isinstance(value, str)
        and len(value) == 7
        and value.startswith("#")
        and all(character in "0123456789abcdefABCDEF" for character in value[1:])
    )


def _declares_path(value: Any, fields: list[str]) -> bool:
    current = value
    for field in fields:
        if not isinstance(current, dict) or field not in current:
            return False
        current = current[field]
    return True


def _style_source_path(
    raw_style: Any,
    defaults_style: Any,
    part_path: str,
    suffix: str,
) -> str:
    fields = [field for field in suffix.removeprefix(".").split(".") if field]
    if fields and _declares_path(raw_style, fields):
        return f"{part_path}.style{suffix}"
    if fields and _declares_path(defaults_style, fields):
        return f"{part_path.rsplit('.', 1)[0]}.part_defaults.style{suffix}"
    return f"{part_path}.style{suffix}"


def _validate_text_style(value: Any, path: str) -> list[ContractIssue]:
    if not isinstance(value, dict):
        return [_issue(path, "text_style must be an object")]
    unknown = sorted(set(value) - _TEXT_STYLE_FIELDS)
    missing = sorted(_TEXT_STYLE_FIELDS - set(value))
    if unknown:
        return [_issue(path, f"unknown text_style fields: {', '.join(unknown)}")]
    if missing:
        return [_issue(path, f"missing text_style fields: {', '.join(missing)}")]
    margins = value["margins"]
    if not isinstance(margins, dict) or set(margins) != _MARGIN_FIELDS:
        return [_issue(f"{path}.margins", "all four margins must be non-negative integer EMU")]
    for side in sorted(_MARGIN_FIELDS):
        if not valid_nonnegative_coordinate32(margins[side]):
            return [_issue(
                f"{path}.margins.{side}",
                "margin must be a non-negative integer EMU no greater than 2147483647",
            )]
    if (
        not isinstance(value["font_name"], str)
        or not value["font_name"]
        or not valid_font_size_pt(value["font_size"])
        or type(value["font_weight"]) is not int
        or not 1 <= value["font_weight"] <= 1000
        or not _valid_rgb(value["color"])
        or type(value["italic"]) is not bool
        or not isinstance(value["alignment"], str)
        or value["alignment"] not in _ALIGNMENTS
        or not isinstance(value["vertical_alignment"], str)
        or value["vertical_alignment"] not in _VERTICAL_ALIGNMENTS
        or type(value["wrap"]) is not bool
    ):
        return [_issue(path, "text_style values are invalid")]
    return []


def _set_font_family(run: Any, family: str) -> None:
    run.font.name = family
    properties = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs", "a:sym"):
        node = properties.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            properties.append(node)
        node.set("typeface", family)


def _apply_text(shape: Any, text: str, style: dict[str, Any]) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.auto_size = MSO_AUTO_SIZE.NONE
    margins = style["margins"]
    frame.margin_left = Emu(margins["left"])
    frame.margin_right = Emu(margins["right"])
    frame.margin_top = Emu(margins["top"])
    frame.margin_bottom = Emu(margins["bottom"])
    frame.word_wrap = style["wrap"]
    frame.vertical_anchor = _VERTICAL_ALIGNMENTS[style["vertical_alignment"]]
    paragraph = frame.paragraphs[0]
    paragraph.alignment = _ALIGNMENTS[style["alignment"]]
    run = paragraph.add_run()
    run.text = text
    _set_font_family(run, style["font_name"])
    run.font.size = Pt(style["font_size"])
    run.font.bold = style["font_weight"] >= 600
    run.font.italic = style["italic"]
    run.font.color.rgb = RGBColor.from_string(style["color"][1:])


class MultipartRenderer:
    supported_fields = frozenset(
        {"rotation", "part_defaults", "parts", "repeat_sequence", "allow_overlap"}
    )
    supported_values: dict[str, frozenset[str]] = {}
    required_fields = frozenset({"part_defaults"})
    capability_ids = frozenset({"multipart.repeat_sequence"})

    def __init__(self, kind: str) -> None:
        self.kind = kind

    def validate_contract(
        self, element: dict[str, Any], context: RenderContext
    ) -> list[ContractIssue]:
        issues = validate_element_contract(element)
        if issues:
            return issues
        element_id = element["element_id"]
        if context.representation_modes.get(element_id) != "composite":
            return [_issue(
                f"elements.{element_id}.representation",
                "multipart renderer requires composite mode",
            )]
        if element["style"].get("rotation", 0) != 0:
            return [_issue(
                f"elements.{element_id}.style.rotation",
                "multipart parent rotation must be zero; declare rotation per part",
            )]
        sequence_name = (
            "parts" if "parts" in element["content"] else "repeat_sequence"
        )
        defaults_style = element["content"]["part_defaults"].get("style", {})
        raw_parts = element["content"][sequence_name]
        for index, part in enumerate(expand_multipart_parts(element)):
            part_path = f"elements.{element_id}.content.{sequence_name}[{index}]"
            raw_part = raw_parts[index]
            raw_style = raw_part.get("style", {})
            style = part.get("style")
            if not isinstance(style, dict):
                return [_issue(f"{part_path}.style", "part style is required")]
            shape_style = _part_shape_style(style)
            shape_issues = SHAPE_RENDERER.validate_contract(
                {
                    "element_id": f"{element_id}:{part.get('part_id', index)}",
                    "style": shape_style,
                    "content": {},
                },
                context,
            )
            if shape_issues:
                synthetic_prefix = (
                    f"elements.{element_id}:{part.get('part_id', index)}.style"
                )
                issue = shape_issues[0]
                issue_path = issue.path
                if issue_path.startswith(synthetic_prefix):
                    suffix = issue_path[len(synthetic_prefix):]
                    issue_path = _style_source_path(
                        raw_style,
                        defaults_style,
                        part_path,
                        suffix,
                    )
                return [
                    ContractIssue(
                        issue.code,
                        issue_path,
                        issue.detail,
                        issue.capability,
                    )
                ]
            content = part.get("content", {})
            has_text = isinstance(content, dict) and "text" in content
            text = content.get("text") if isinstance(content, dict) else None
            if has_text:
                if not isinstance(text, str):
                    return [_issue(f"{part_path}.content.text", "part text must be a string")]
                text_style_path = f"{part_path}.style.text_style"
                text_issues = _validate_text_style(
                    style.get("text_style"), text_style_path
                )
                if text_issues:
                    issue = text_issues[0]
                    suffix = issue.path[len(text_style_path):]
                    source = _style_source_path(
                        raw_style,
                        defaults_style,
                        part_path,
                        f".text_style{suffix}",
                    )
                    return [
                        ContractIssue(
                            issue.code,
                            source,
                            issue.detail,
                            issue.capability,
                        )
                    ]
        return []

    def render(self, element: dict[str, Any], context: RenderContext) -> None:
        element_id = element["element_id"]
        actual_boxes: list[list[int]] = []
        for part in expand_multipart_parts(element):
            style = _part_shape_style(part["style"])
            x, y, width, height = part["slide_bbox"]
            shape = context.slide.shapes.add_shape(
                _SHAPE_TYPES[style["shape_type"]],
                Emu(x),
                Emu(y),
                Emu(width),
                Emu(height),
            )
            _apply_fill(shape, style.get("fill"))
            if "line" in style:
                apply_stroke(shape, style["line"])
            else:
                shape.line.fill.background()
            _apply_shadow(shape, style.get("effects"))
            if "adjustments" in style:
                set_preset_shape_adjustments(
                    shape,
                    style["adjustments"],
                    f"elements.{element_id}.parts.{part['part_id']}.style.adjustments",
                )
            set_shape_flips(
                shape,
                style.get("flip_horizontal", False),
                style.get("flip_vertical", False),
                f"elements.{element_id}.parts.{part['part_id']}.style",
            )
            shape.rotation = style.get("rotation", 0)
            content = part.get("content", {})
            has_text = isinstance(content, dict) and "text" in content
            text = content.get("text") if isinstance(content, dict) else None
            fonts: tuple[str, ...] = ()
            if has_text:
                if not isinstance(text, str):
                    raise ToolError(
                        "UNSUPPORTED_CAPABILITY",
                        f"elements.{element_id}.parts.{part['part_id']}.content.text",
                        "part text must be a string",
                    )
                _apply_text(shape, text, part["style"]["text_style"])
                fonts = (part["style"]["text_style"]["font_name"],)
            context.registry.register(
                element_id,
                shape,
                "sp",
                semantic_kind=self.kind,
                selected_mode=context.representation_modes[element_id],
                part_id=part["part_id"],
                text_summary=text,
                font_declarations=fonts,
            )
            actual_boxes.append(
                [int(shape.left), int(shape.top), int(shape.width), int(shape.height)]
            )
        parent = validate_bbox(
            element["slide_bbox"], f"elements.{element_id}.slide_bbox"
        )
        if bbox_union(actual_boxes) != parent:
            raise ToolError(
                "BUILD_OUTPUT_INCOMPLETE",
                f"elements.{element_id}",
                "rendered multipart union does not equal the parent bbox",
            )


MATRIX_RENDERER = MultipartRenderer("matrix")
STATUS_RENDERER = MultipartRenderer("status")
register_renderer("matrix", MATRIX_RENDERER)
register_renderer("status", STATUS_RENDERER)
